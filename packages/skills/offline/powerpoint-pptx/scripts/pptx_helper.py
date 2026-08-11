#!/usr/bin/env python3
"""Deterministic, non-destructive PPTX inspection and slide append commands."""

from __future__ import annotations

import argparse
from collections import Counter
import hashlib
import json
import os
from pathlib import Path
import tempfile
from typing import Callable

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


def file_digest(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as source:
        for block in iter(lambda: source.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def input_path(raw: str) -> Path:
    path = Path(raw).expanduser().resolve(strict=True)
    if not path.is_file() or path.suffix.lower() != ".pptx":
        raise ValueError(f"input must be an existing .pptx file: {path}")
    return path


def output_path(raw: str, source: Path) -> Path:
    path = Path(raw).expanduser().resolve()
    if path.suffix.lower() != ".pptx":
        raise ValueError(f"output must use the .pptx extension: {path}")
    if path == source:
        raise ValueError("output must be different from input")
    if path.exists():
        raise FileExistsError(f"output already exists: {path}")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


def slide_texts(presentation) -> list[list[str]]:
    return [
        [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]
        for slide in presentation.slides
    ]


def save_verified(
    presentation,
    source: Path,
    output: Path,
    verify: Callable[[object], None],
) -> None:
    source_before = file_digest(source)
    temporary: Path | None = None
    linked = False
    try:
        with tempfile.NamedTemporaryFile(
            prefix=f".{output.name}.", suffix=".tmp", dir=output.parent, delete=False
        ) as handle:
            temporary = Path(handle.name)
        presentation.save(temporary)
        verify(Presentation(temporary))
        os.link(temporary, output)
        linked = True
        verify(Presentation(output))
        if file_digest(source) != source_before:
            raise RuntimeError("source PPTX changed during editing")
    except Exception:
        if linked:
            output.unlink(missing_ok=True)
        raise
    finally:
        if temporary is not None:
            temporary.unlink(missing_ok=True)


def inspect_presentation(args: argparse.Namespace) -> dict[str, object]:
    source = input_path(args.input)
    presentation = Presentation(source)
    slides = [
        {"index": index, "texts": texts}
        for index, texts in enumerate(slide_texts(presentation), start=1)
    ]
    return {
        "operation": "inspect",
        "input": str(source),
        "slide_count": len(slides),
        "slides": slides,
        "verified": True,
    }


def add_text_box(slide, presentation, *, top_fraction: int, height_fraction: int):
    left = presentation.slide_width // 10
    width = presentation.slide_width * 8 // 10
    top = presentation.slide_height * top_fraction // 10
    height = presentation.slide_height * height_fraction // 10
    return slide.shapes.add_textbox(left, top, width, height)


def append_slide(args: argparse.Namespace) -> dict[str, object]:
    source = input_path(args.input)
    output = output_path(args.output, source)
    presentation = Presentation(source)
    before = slide_texts(presentation)
    if len(presentation.slide_layouts) == 0:
        raise ValueError("input PPTX has no slide layout available for a new slide")

    layout_index = 1 if len(presentation.slide_layouts) > 1 else 0
    slide = presentation.slides.add_slide(presentation.slide_layouts[layout_index])
    title_shape = slide.shapes.title
    if title_shape is None:
        title_shape = add_text_box(slide, presentation, top_fraction=1, height_fraction=1)
    title_shape.text = args.title

    body_shape = next(
        (
            shape
            for shape in slide.placeholders
            if shape.placeholder_format.type
            in {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE}
        ),
        None,
    )
    if body_shape is None:
        body_shape = add_text_box(slide, presentation, top_fraction=3, height_fraction=4)
    body_shape.text = args.body

    def verify(reopened) -> None:
        texts = slide_texts(reopened)
        if len(texts) != len(before) + 1 or texts[: len(before)] != before:
            raise RuntimeError("existing slides changed or slide count differs after reopen")
        actual = Counter(texts[-1])
        required = Counter((args.title, args.body))
        if any(actual[text] < count for text, count in required.items()):
            raise RuntimeError("new slide content is missing after reopen")

    save_verified(presentation, source, output, verify)
    return {
        "operation": "append-slide",
        "input": str(source),
        "output": str(output),
        "slide_count": len(before) + 1,
        "title": args.title,
        "body": args.body,
        "verified": True,
    }


def parser() -> argparse.ArgumentParser:
    root = argparse.ArgumentParser(description=__doc__)
    commands = root.add_subparsers(dest="command", required=True)

    inspect = commands.add_parser("inspect", help="read slide text")
    inspect.add_argument("--input", required=True)
    inspect.set_defaults(run=inspect_presentation)

    append = commands.add_parser("append-slide", help="append one title-and-body slide")
    append.add_argument("--input", required=True)
    append.add_argument("--output", required=True)
    append.add_argument("--title", required=True)
    append.add_argument("--body", required=True)
    append.set_defaults(run=append_slide)
    return root


def main() -> int:
    args = parser().parse_args()
    try:
        result = args.run(args)
    except Exception as error:
        print(json.dumps({"error": str(error), "verified": False}, ensure_ascii=False))
        return 1
    print(json.dumps(result, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
