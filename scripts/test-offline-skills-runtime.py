#!/usr/bin/env python3
"""Exercise all offline Skill bootstraps with one native, pre-downloaded wheelhouse."""

from __future__ import annotations

import argparse
import hashlib
import os
from pathlib import Path
import shutil
import subprocess
import sys
import tempfile


def run(command: list[str], *, env: dict[str, str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        command,
        check=False,
        env=env,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        timeout=360,
    )


def checked(command: list[str], *, env: dict[str, str], operation: str) -> None:
    result = run(command, env=env)
    if result.returncode != 0:
        raise RuntimeError(f"{operation} failed ({result.returncode}):\n{result.stdout}")


def digest(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def environment_python(agent: Path, skill_name: str) -> Path:
    environment = agent / "shared_env" / skill_name
    if os.name == "nt":
        return environment / "Scripts" / "python.exe"
    return environment / "bin" / "python"


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--wheel-dir", type=Path, required=True)
    args = parser.parse_args()

    wheel_dir = args.wheel_dir.resolve(strict=True)
    wheels = list(wheel_dir.glob("*.whl"))
    if len(wheels) != 7:
        raise RuntimeError(f"expected two native and five universal wheels, found {len(wheels)}")

    repository = Path(__file__).resolve().parent.parent
    source_root = repository / "packages" / "skills" / "offline"
    with tempfile.TemporaryDirectory(prefix="penguin-offline-skills-native-") as temporary:
        root = Path(temporary)
        agent = root / "project" / "agents" / "default_agent"
        skill_root = agent / "agent_state" / "skills"
        for skill_name in ("word-docx", "powerpoint-pptx", "pdf-tools"):
            shutil.copytree(source_root / skill_name, skill_root / skill_name)

        offline_root = root / "offline"
        (offline_root / "_shared").mkdir(parents=True)
        shutil.copy2(
            source_root / "_shared" / "bootstrap_runtime.py",
            offline_root / "_shared" / "bootstrap_runtime.py",
        )
        (offline_root / "wheels").mkdir()
        for wheel in wheels:
            shutil.copy2(wheel, offline_root / "wheels" / wheel.name)

        environment = os.environ.copy()
        for key in tuple(environment):
            if key.startswith("PIP_") or key in {"PYTHONHOME", "PYTHONPATH"}:
                environment.pop(key)
        escaped = root / "escaped-pip-target"
        environment.update(
            {
                "PENGUIN_OFFLINE_ROOT": str(offline_root),
                "PIP_TARGET": str(escaped),
                "PIP_USER": "1",
                "PYTHONHOME": str(root / "invalid-python-home"),
                "PYTHONPATH": str(root / "invalid-python-path"),
            }
        )

        bootstraps = {
            name: skill_root / name / "scripts" / "bootstrap.py"
            for name in ("word-docx", "powerpoint-pptx", "pdf-tools")
        }
        # Two simultaneous first calls exercise the cross-platform lock implementation.
        processes = [
            subprocess.Popen(
                [sys.executable, "-I", str(bootstraps["word-docx"]), "--help"],
                env=environment,
                stdout=subprocess.PIPE,
                stderr=subprocess.STDOUT,
                text=True,
            )
            for _ in range(2)
        ]
        for process in processes:
            output, _ = process.communicate(timeout=360)
            if process.returncode != 0:
                raise RuntimeError(f"concurrent DOCX bootstrap failed ({process.returncode}):\n{output}")
        for name in ("powerpoint-pptx", "pdf-tools"):
            checked(
                [sys.executable, "-I", str(bootstraps[name]), "--help"],
                env=environment,
                operation=f"{name} bootstrap",
            )
        if escaped.exists():
            raise RuntimeError("pip environment escaped the Agent-owned virtual environments")

        word_python = environment_python(agent, "word-docx")
        pptx_python = environment_python(agent, "powerpoint-pptx")
        pdf_python = environment_python(agent, "pdf-tools")
        for python, imports in (
            (word_python, "import docx, lxml"),
            (pptx_python, "import pptx, lxml, PIL, xlsxwriter"),
            (pdf_python, "import pypdf"),
        ):
            checked([str(python), "-I", "-c", imports], env=environment, operation="import probe")

        # A matching marker must not make a damaged cached environment permanent. Remove one
        # installed package, then prove the next bootstrap call rebuilds the environment offline.
        location = run(
            [
                str(word_python),
                "-I",
                "-c",
                "from pathlib import Path; import docx; print(Path(docx.__file__).resolve().parent)",
            ],
            env=environment,
        )
        if location.returncode != 0:
            raise RuntimeError(f"DOCX package lookup failed:\n{location.stdout}")
        shutil.rmtree(Path(location.stdout.strip()))
        checked(
            [sys.executable, "-I", str(bootstraps["word-docx"]), "--help"],
            env=environment,
            operation="damaged DOCX environment recovery",
        )
        checked(
            [str(word_python), "-I", "-c", "import docx, lxml"],
            env=environment,
            operation="recovered DOCX import probe",
        )

        docx_input = root / "input-without-heading.docx"
        docx_output = root / "output.docx"
        checked(
            [
                str(word_python),
                "-I",
                "-c",
                (
                    "from docx import Document; "
                    "d=Document(); h=d.styles['Heading 1']; d.styles.element.remove(h._element); "
                    "d.add_paragraph('Existing paragraph'); d.save(r'%s')"
                )
                % docx_input,
            ],
            env=environment,
            operation="DOCX fixture creation",
        )
        docx_hash = digest(docx_input)
        checked(
            [
                sys.executable,
                "-I",
                str(bootstraps["word-docx"]),
                "append",
                "--input",
                str(docx_input),
                "--output",
                str(docx_output),
                "--heading",
                "Native offline heading",
                "--paragraph",
                "Native offline paragraph",
            ],
            env=environment,
            operation="DOCX append",
        )
        if digest(docx_input) != docx_hash:
            raise RuntimeError("source DOCX changed")
        checked(
            [
                str(word_python),
                "-I",
                "-c",
                (
                    "from docx import Document; d=Document(r'%s'); texts=[p.text for p in d.paragraphs]; "
                    "assert 'Existing paragraph' in texts; assert 'Native offline heading' in texts; "
                    "assert 'Native offline paragraph' in texts; "
                    "p=next(p for p in d.paragraphs if p.text=='Native offline heading'); "
                    "assert any(r.bold is True for r in p.runs)"
                )
                % docx_output,
            ],
            env=environment,
            operation="DOCX reopen verification",
        )

        # Assigning Run.text rebuilds the run XML. Refuse replacement when a matching run also
        # carries an inline image instead of silently deleting the drawing and reporting success.
        mixed_input = root / "input-mixed-run.docx"
        mixed_output = root / "output-mixed-run.docx"
        checked(
            [
                str(word_python),
                "-I",
                "-c",
                (
                    "import base64; from io import BytesIO; from docx import Document; "
                    "d=Document(); r=d.add_paragraph().add_run('replace me '); "
                    "r.add_picture(BytesIO(base64.b64decode("
                    "'iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII='"
                    "))); r.add_text(' safely'); d.save(r'%s')"
                )
                % mixed_input,
            ],
            env=environment,
            operation="mixed-run DOCX fixture creation",
        )
        mixed_hash = digest(mixed_input)
        rejected = run(
            [
                sys.executable,
                "-I",
                str(bootstraps["word-docx"]),
                "replace",
                "--input",
                str(mixed_input),
                "--output",
                str(mixed_output),
                "--old",
                "replace me",
                "--new",
                "replaced",
            ],
            env=environment,
        )
        if rejected.returncode == 0 or mixed_output.exists():
            raise RuntimeError("DOCX replacement did not reject a matching run with an image")
        if digest(mixed_input) != mixed_hash:
            raise RuntimeError("mixed-run source DOCX changed")
        checked(
            [
                str(word_python),
                "-I",
                "-c",
                "from docx import Document; assert len(Document(r'%s').inline_shapes)==1"
                % mixed_input,
            ],
            env=environment,
            operation="mixed-run DOCX preservation check",
        )

        pptx_input = root / "input.pptx"
        pptx_output = root / "output.pptx"
        checked(
            [
                str(pptx_python),
                "-I",
                "-c",
                (
                    "from pptx import Presentation\n"
                    "presentation = Presentation()\n"
                    "for index in range(1, 4):\n"
                    "    slide = presentation.slides.add_slide(presentation.slide_layouts[1])\n"
                    "    slide.shapes.title.text = f'Slide {index}'\n"
                    "    slide.placeholders[1].text = f'Body {index}'\n"
                    "presentation.save(r'%s')\n"
                )
                % pptx_input,
            ],
            env=environment,
            operation="PPTX fixture creation",
        )
        pptx_hash = digest(pptx_input)
        checked(
            [
                sys.executable,
                "-I",
                str(bootstraps["powerpoint-pptx"]),
                "append-slide",
                "--input",
                str(pptx_input),
                "--output",
                str(pptx_output),
                "--title",
                "Slide 4",
                "--body",
                "Offline PPTX body",
            ],
            env=environment,
            operation="PPTX append",
        )
        if digest(pptx_input) != pptx_hash:
            raise RuntimeError("source PPTX changed")
        checked(
            [
                str(pptx_python),
                "-I",
                "-c",
                (
                    "from pptx import Presentation; p=Presentation(r'%s'); assert len(p.slides)==4; "
                    "texts=[s.text for s in p.slides[-1].shapes if getattr(s,'has_text_frame',False)]; "
                    "assert 'Slide 4' in texts and 'Offline PPTX body' in texts"
                )
                % pptx_output,
            ],
            env=environment,
            operation="PPTX reopen verification",
        )

        same_text_output = root / "output-same-title-body.pptx"
        checked(
            [
                sys.executable,
                "-I",
                str(bootstraps["powerpoint-pptx"]),
                "append-slide",
                "--input",
                str(pptx_input),
                "--output",
                str(same_text_output),
                "--title",
                "Same text",
                "--body",
                "Same text",
            ],
            env=environment,
            operation="PPTX identical title/body append",
        )
        checked(
            [
                str(pptx_python),
                "-I",
                "-c",
                (
                    "from pptx import Presentation; p=Presentation(r'%s'); "
                    "texts=[s.text for s in p.slides[-1].shapes if getattr(s,'has_text_frame',False)]; "
                    "assert texts.count('Same text')>=2"
                )
                % same_text_output,
            ],
            env=environment,
            operation="PPTX identical title/body verification",
        )

        pdf_main = root / "main.pdf"
        pdf_appendix = root / "appendix.pdf"
        pdf_output = root / "output.pdf"
        checked(
            [
                str(pdf_python),
                "-I",
                "-c",
                (
                    "from pypdf import PdfWriter; "
                    "a=PdfWriter(); a.add_blank_page(200,500); a.add_blank_page(300,500); "
                    "f=open(r'%s','wb'); a.write(f); f.close(); "
                    "b=PdfWriter(); b.add_blank_page(400,500); "
                    "f=open(r'%s','wb'); b.write(f); f.close()"
                )
                % (pdf_main, pdf_appendix),
            ],
            env=environment,
            operation="PDF fixture creation",
        )
        pdf_hashes = (digest(pdf_main), digest(pdf_appendix))
        checked(
            [
                sys.executable,
                "-I",
                str(bootstraps["pdf-tools"]),
                "merge",
                "--input",
                str(pdf_main),
                "--input",
                str(pdf_appendix),
                "--output",
                str(pdf_output),
            ],
            env=environment,
            operation="PDF merge",
        )
        if pdf_hashes != (digest(pdf_main), digest(pdf_appendix)):
            raise RuntimeError("source PDF changed")
        checked(
            [
                str(pdf_python),
                "-I",
                "-c",
                (
                    "from pypdf import PdfReader; p=PdfReader(r'%s'); "
                    "assert [float(x.mediabox.width) for x in p.pages]==[200.0,300.0,400.0]"
                )
                % pdf_output,
            ],
            env=environment,
            operation="PDF reopen verification",
        )

    platform = os.uname().machine if hasattr(os, "uname") else os.environ.get(
        "PROCESSOR_ARCHITECTURE", "unknown"
    )
    print(f"offline Skill native runtime acceptance passed: {sys.platform}/{platform}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
